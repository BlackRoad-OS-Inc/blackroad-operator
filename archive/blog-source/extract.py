#!/usr/bin/env python3
"""Extract text from all .pptx files into individual .md files."""
import os
import sys
from pptx import Presentation

INPUT_DIRS = [
    os.path.expanduser("~/blog-source/pptx/digital-comm"),
    os.path.expanduser("~/blog-source/pptx/strategic-comm"),
]
OUTPUT_DIR = os.path.expanduser("~/blog-source/text")

for input_dir in INPUT_DIRS:
    category = os.path.basename(input_dir)
    for fname in sorted(os.listdir(input_dir)):
        if not fname.endswith(".pptx"):
            continue
        fpath = os.path.join(input_dir, fname)
        prs = Presentation(fpath)
        
        lines = [f"# {fname.replace('.pptx', '')}\n"]
        lines.append(f"**Source:** {category}\n\n---\n")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip(" |"):
                            slide_texts.append(row_text)
            
            if slide_texts:
                lines.append(f"\n## Slide {slide_num}\n")
                for t in slide_texts:
                    lines.append(f"{t}\n")
        
        out_name = fname.replace(".pptx", ".md").replace(" ", "-")
        out_path = os.path.join(OUTPUT_DIR, f"{category}--{out_name}")
        with open(out_path, "w") as f:
            f.write("\n".join(lines))
        print(f"  Extracted: {out_name} ({len(prs.slides)} slides)")

print(f"\nDone! Files in {OUTPUT_DIR}")

# Extract blackroad drive pptx
INPUT_DIRS2 = [os.path.expanduser("~/blog-source/pptx/blackroad")]
for input_dir in INPUT_DIRS2:
    category = "blackroad"
    for fname in sorted(os.listdir(input_dir)):
        if not fname.endswith(".pptx"):
            continue
        fpath = os.path.join(input_dir, fname)
        prs = Presentation(fpath)
        lines = [f"# {fname.replace('.pptx', '')}\n"]
        lines.append(f"**Source:** {category}\n\n---\n")
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip(" |"):
                            slide_texts.append(row_text)
            if slide_texts:
                lines.append(f"\n## Slide {slide_num}\n")
                for t in slide_texts:
                    lines.append(f"{t}\n")
        out_name = fname.replace(".pptx", ".md").replace(" ", "-")
        out_path = os.path.join(OUTPUT_DIR, f"{category}--{out_name}")
        with open(out_path, "w") as f:
            f.write("\n".join(lines))
        print(f"  Extracted: {out_name} ({len(prs.slides)} slides)")
